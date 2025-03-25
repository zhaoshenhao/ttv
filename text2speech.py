import os
import yaml
from pptx import Presentation
import subprocess
import re
from f5_tts_api import F5TTS

class Text2Speech:
    def __init__(self, ppt_file, lang="zh", config_file="config.yaml"):
        self.ppt_file = ppt_file
        self.lang = lang.lower()
        self.prs = Presentation(ppt_file)
        
        with open(config_file, 'r', encoding='utf-8') as f:
            config = yaml.safe_load(f)

        self.ref_zh_audio = config.get('ref_zh_audio', 'tts/F5TTS_zh.wav')
        self.ref_zh_text = config.get('ref_zh_text', 'tts/F5TTS_zh.txt')
        self.ref_en_audio = config.get('ref_en_audio', 'tts/F5TTS_en.wav')
        self.ref_en_text = config.get('ref_en_text', 'tts/F5TTS_en.txt')
        self.model = config.get('model', 'F5TTS_v1_Base')
        self.speed = config.get('speed', 1.0)
        self.vocoder_name = config.get('vocoder_name', 'vocos')
        self.target_rms = config.get('target_rms', 0.1)
        self.remove_silence = config.get('remove_silence', True)
        self.audio_dir = config.get('output_dir', 'audio')        
        
        if not os.path.exists(self.audio_dir):
            os.makedirs(self.audio_dir)

        if self.lang == "zh":
            self.ref_audio = self.ref_zh_audio
            self.ref_text = self.ref_zh_text
        elif self.lang == "en":
            self.ref_audio = self.ref_en_audio
            self.ref_text = self.ref_en_text
        else:
            raise ValueError("Language must be 'zh' or 'en'")

    def split_sentence(self, sentence):
        punctuation = r'[.,!?;:。，！？；：]'
        parts = re.split(f'({punctuation})', sentence)
        result = []
        for i in range(0, len(parts), 2):
            if i + 1 < len(parts):
                result.append(parts[i] + parts[i + 1])
            else:
                result.append(parts[i])
        return result

    def split_text(self, text):
        lines = []
        paragraphs = text.splitlines()
        for para in paragraphs:
            if not para.strip():
                continue
            result = self.split_sentence(para)
            lines.extend(result)
        return lines

    def process_uppercase(self, text):
        words = text.split()
        processed_words = []
        
        for word in words:
            if word.isupper() and word.isalpha():
                processed_word = ".".join(word) + "."
                processed_words.append(processed_word)
            else:
                processed_words.append(word)
        
        return " ".join(processed_words)
    
    def remove_punctuation(self, text):
        punctuation = r'[.,!?;:。，！？；：]'
        return re.sub(punctuation, '', text)

    def generate_text_files(self):
        for i, slide in enumerate(self.prs.slides):
            text = slide.notes_slide.notes_text_frame.text.strip()
            if not text:
                print(f"No notes found for slide {i}, skipping text file generation")
                continue
            
            lines = self.split_text(text)
            processed_lines = [self.process_uppercase(line) for line in lines]
            
            for idx, l in enumerate(processed_lines):
                if not self.remove_punctuation(l.strip()):
                    continue
                ii = f"{i:03}"
                idxs = f"{idx:03}"
                txt_file = os.path.join(self.audio_dir, f"slide-{ii}-{idxs}.txt")
                with open(txt_file, 'w', encoding='utf-8') as f:
                    f.write(l)
                print(f"Generated text file for slide {i}: {txt_file}")

    def generate_audio(self):
        txt_files = []
        for file in os.listdir(self.audio_dir):
            if file.endswith(".txt"):
                txt_files.append(file)
        
        f5tts = F5TTS()
        for file in txt_files:
            txt_file = os.path.join(self.audio_dir, file)
            audio_file = txt_file[:-3] + 'wav'
            with open(txt_file, "r", encoding="utf-8") as f:
                s = f.read()
            try:
                print(f"Text file: {txt_file}, Audio file: {audio_file}")
                f5tts.infer(
                    ref_file = self.ref_audio,
                    ref_text = self.ref_text,
                    gen_text = s,
                    target_rms = self.target_rms,
                    speed = self.speed,
                    remove_silence= self.remove_silence,
                    file_wave= audio_file,
                )
            except Exception as e:
                print(f"Error generating audio: {e.stderr}")

    def convert(self):
        """主流程：生成文本文件和音频"""
        self.generate_text_files()
        self.generate_audio()