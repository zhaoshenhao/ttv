import os
from docx import Document
from pptx import Presentation
from docx.oxml.ns import qn
import io

class Word2PPTX:
    def __init__(self, input_doc, output_ppt, template_ppt, max_leaf_count=8):
        self.input_doc = input_doc
        self.output_ppt = output_ppt
        self.template_ppt = template_ppt
        self.max_leaf_count = max_leaf_count
        self.doc = Document(input_doc)
        self.prs = Presentation(template_ppt)
        if not template_ppt:
            raise ValueError("A template PPT file must be provided")

    def count_leaf_headings(self, start_idx, end_idx):
        """统计指定范围内叶子标题数量（2-4级，不含子标题）"""
        leaf_count = 0
        i = start_idx
        while i < end_idx:
            para = self.doc.paragraphs[i]
            style = para.style.name
            if "Heading" in style:
                level = int(style.split()[-1]) if style.split()[-1].isdigit() else 0
                if 2 <= level <= 4:
                    is_leaf = True
                    for j in range(i + 1, end_idx):
                        next_style = self.doc.paragraphs[j].style.name
                        if "Heading" in next_style and int(next_style.split()[-1]) > level:
                            is_leaf = False
                            break
                    if is_leaf:
                        leaf_count += 1
            i += 1
        return leaf_count

    def extract_images(self, start_idx, end_idx):
        """提取指定范围内的图片，包括内联和浮动图片"""
        images = []
        print(f"Checking images between paragraphs {start_idx} and {end_idx}")
        
        for idx, shape in enumerate(self.doc.inline_shapes):
            if shape.type == 3:  # Picture
                shape_element = shape._inline
                for i in range(start_idx, end_idx):
                    para = self.doc.paragraphs[i]
                    for run_idx, run in enumerate(para.runs):
                        if shape_element in run._element.getparent().getchildren():
                            image_rid = shape._inline.graphic.graphicData.pic.blipFill.blip.embed
                            try:
                                image_blob = self.doc.part.related_parts[image_rid].blob
                                images.append((i, image_blob))
                                print(f"    Extracted inline image at paragraph {i} (RID: {image_rid})")
                                break
                            except KeyError:
                                print(f"    Warning: Inline image RID {image_rid} not found")
                                break
        
        for rel in self.doc.part.rels.values():
            if "image" in rel.target_ref:
                for i in range(start_idx, end_idx):
                    para = self.doc.paragraphs[i]
                    for run in para.runs:
                        drawing_elements = run._element.findall(qn('w:drawing'))
                        for drawing in drawing_elements:
                            blip = drawing.find('.//' + qn('a:blip'))
                            if blip is not None and blip.embed == rel.rId:
                                try:
                                    image_blob = self.doc.part.related_parts[rel.rId].blob
                                    images.append((i, image_blob))
                                    print(f"  Extracted floating image at paragraph {i} (RID: {rel.rId})")
                                    break
                                except KeyError:
                                    print(f"  Warning: Floating image RID {rel.rId} not found")
                                break
        
        return sorted(images, key=lambda x: x[0])

    def add_slide(self, layout_idx, title, subheadings=None, notes=""):
        """添加幻灯片，不处理图片缩放"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])
        
        slide.shapes.title.text = title
        
        if subheadings and len(subheadings) > 0:
            content_shape = None
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.idx == 1:
                    content_shape = shape
                    break
            if content_shape:
                text_frame = content_shape.text_frame
                text_frame.clear()
                for i, subheading in enumerate(subheadings):
                    p = text_frame.add_paragraph()
                    p.text = subheading
                    if i > 0:
                        p.level = 0
        
        if notes:
            slide.notes_slide.notes_text_frame.text = notes.strip()
        
        print(f"Added slide {len(self.prs.slides)-1}: Title='{title}', Subheadings={len(subheadings) if subheadings else 0}, Notes='{notes.strip()[:50] if notes else ''}...'")
        return slide

    def add_image_slide(self, layout_idx, title, image_blob, notes=""):
        """添加单独的图片幻灯片，不缩放"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])
        
        slide.shapes.title.text = title
        
        pic = slide.shapes.add_picture(io.BytesIO(image_blob), left=0, top=0)
        print(f"Added image slide {len(self.prs.slides)-1}: Title='{title}', Image size={pic.width}x{pic.height}, Notes='{notes.strip()[:50] if notes else ''}...'")
        
        if notes:
            slide.notes_slide.notes_text_frame.text = notes.strip()
        
        return slide

    def convert(self):
        """执行Word到PPT转换"""
        # 第一页：文章标题
        title = None
        title_idx = -1
        for i, para in enumerate(self.doc.paragraphs):
            if para.style.name == "Title" and para.text.strip():
                title = para.text.strip()
                title_idx = i
                break
        if not title and self.doc.paragraphs and self.doc.paragraphs[0].text.strip():
            title = self.doc.paragraphs[0].text.strip()
            title_idx = 0
        if title:
            if self.prs.slides:
                self.prs.slides[0].shapes.title.text = title
            else:
                self.add_slide(0, title)
        else:
            if self.prs.slides:
                self.prs.slides[0].shapes.title.text = "Untitled Document"
            else:
                self.add_slide(0, "Untitled Document")
            title_idx = -1

        # 第二页：Agenda（仅Heading 1）
        agenda_notes = ""
        first_h1_idx = -1
        for i, para in enumerate(self.doc.paragraphs):
            if para.style.name == "Heading 1" and para.text.strip():
                first_h1_idx = i
                break
        if title_idx != -1 and first_h1_idx != -1 and first_h1_idx > title_idx + 1:
            for i in range(title_idx + 1, first_h1_idx):
                text = self.doc.paragraphs[i].text.strip()
                if text:
                    agenda_notes += text + "\n"

        toc_subheadings = [para.text.strip() for para in self.doc.paragraphs if para.style.name == "Heading 1" and para.text.strip()]
        if toc_subheadings:
            if len(self.prs.slides) > 1:
                slide = self.prs.slides[1]
                slide.shapes.title.text = "Agenda"
                content_shape = None
                for shape in slide.shapes:
                    if shape.is_placeholder and shape.placeholder_format.idx == 1:
                        content_shape = shape
                        break
                if content_shape:
                    text_frame = content_shape.text_frame
                    text_frame.clear()
                    for i, subheading in enumerate(toc_subheadings):
                        p = text_frame.add_paragraph()
                        p.text = subheading
                        if i > 0:
                            p.level = 0
                if agenda_notes:
                    slide.notes_slide.notes_text_frame.text = agenda_notes.strip()
            else:
                self.add_slide(1, "Agenda", toc_subheadings, agenda_notes)

        # 找到所有 Heading 1 的范围
        sections = []
        start_idx = 0
        for i, para in enumerate(self.doc.paragraphs):
            if "Heading 1" in para.style.name and i > 0:
                sections.append((start_idx, i))
                start_idx = i
        sections.append((start_idx, len(self.doc.paragraphs)))

        # 处理每个大章节
        slide_idx = 2
        for start_idx, end_idx in sections[1:]:
            section_title = self.doc.paragraphs[start_idx].text.strip()
            subheadings = [para.text.strip() for i, para in enumerate(self.doc.paragraphs[start_idx+1:end_idx]) 
                          if "Heading" in para.style.name and "Heading 1" not in para.style.name and para.text.strip()]
            images = self.extract_images(start_idx, end_idx)

            if not images:  # 无图片
                notes = ""
                for i in range(start_idx + 1, end_idx):
                    text = self.doc.paragraphs[i].text.strip()
                    if text:
                        notes += text + "\n"
                if len(self.prs.slides) > slide_idx:
                    slide = self.prs.slides[slide_idx]
                    slide.shapes.title.text = section_title
                    content_shape = None
                    for shape in slide.shapes:
                        if shape.is_placeholder and shape.placeholder_format.idx == 1:
                            content_shape = shape
                            break
                    if content_shape and subheadings:
                        text_frame = content_shape.text_frame
                        text_frame.clear()
                        for i, subheading in enumerate(subheadings):
                            p = text_frame.add_paragraph()
                            p.text = subheading
                            if i > 0:
                                p.level = 0
                    if notes:
                        slide.notes_slide.notes_text_frame.text = notes.strip()
                else:
                    self.add_slide(1, section_title, subheadings, notes)
                slide_idx += 1
            else:  # 有图片
                current_notes = ""
                last_content_idx = start_idx

                # 添加第一个 subheadings slide（到第一张图片之前）
                for i in range(start_idx + 1, images[0][0]):
                    text = self.doc.paragraphs[i].text.strip()
                    if text:
                        current_notes += text + "\n"
                if len(self.prs.slides) > slide_idx:
                    slide = self.prs.slides[slide_idx]
                    slide.shapes.title.text = section_title
                    content_shape = None
                    for shape in slide.shapes:
                        if shape.is_placeholder and shape.placeholder_format.idx == 1:
                            content_shape = shape
                            break
                    if content_shape and subheadings:
                        text_frame = content_shape.text_frame
                        text_frame.clear()
                        for i, subheading in enumerate(subheadings):
                            p = text_frame.add_paragraph()
                            p.text = subheading
                            if i > 0:
                                p.level = 0
                    if current_notes:
                        slide.notes_slide.notes_text_frame.text = current_notes.strip()
                else:
                    self.add_slide(1, section_title, subheadings, current_notes)
                slide_idx += 1

                # 处理图片和后续内容
                for img_idx, (img_para_idx, image_blob) in enumerate(images):
                    # 图片后的文字，直到下一个 Heading 或下一个图片/章节结束
                    image_notes = ""
                    next_stop_idx = end_idx if img_idx == len(images) - 1 else images[img_idx + 1][0]
                    for j in range(img_para_idx + 1, next_stop_idx):
                        text = self.doc.paragraphs[j].text.strip()
                        if text and "Heading" in self.doc.paragraphs[j].style.name:
                            break
                        if text:
                            image_notes += text + "\n"
                    self.add_image_slide(1, section_title, image_blob, image_notes)
                    slide_idx += 1

                    # 图片后的 subheadings slide 内容，从第一个 Heading 开始
                    current_notes = ""
                    next_heading_idx = next_stop_idx
                    for j in range(img_para_idx + 1, next_stop_idx):
                        if "Heading" in self.doc.paragraphs[j].style.name:
                            next_heading_idx = j
                            break
                    if next_heading_idx < next_stop_idx:
                        text = self.doc.paragraphs[next_heading_idx].text.strip()
                        if text:
                            current_notes = text + "\n"
                        for j in range(next_heading_idx + 1, next_stop_idx):
                            text = self.doc.paragraphs[j].text.strip()
                            if text:
                                current_notes += text + "\n"
                    
                    # 只在有内容时添加 subheadings slide
                    if current_notes or (img_idx < len(images) - 1):  # 如果有内容或还有后续图片
                        if len(self.prs.slides) > slide_idx:
                            slide = self.prs.slides[slide_idx]
                            slide.shapes.title.text = section_title
                            content_shape = None
                            for shape in slide.shapes:
                                if shape.is_placeholder and shape.placeholder_format.idx == 1:
                                    content_shape = shape
                                    break
                            if content_shape and subheadings:
                                text_frame = content_shape.text_frame
                                text_frame.clear()
                                for i, subheading in enumerate(subheadings):
                                    p = text_frame.add_paragraph()
                                    p.text = subheading
                                    if i > 0:
                                        p.level = 0
                            if current_notes:
                                slide.notes_slide.notes_text_frame.text = current_notes.strip()
                        else:
                            self.add_slide(1, section_title, subheadings, current_notes)
                        slide_idx += 1

        self.prs.save(self.output_ppt)
        print(f"PPT saved as {self.output_ppt} with {len(self.prs.slides)} slides")
        return self.output_ppt