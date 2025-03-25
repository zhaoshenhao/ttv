import os
import yaml
from pptx import Presentation
from moviepy import ImageClip, AudioFileClip, concatenate_videoclips, concatenate_audioclips
import win32com.client  # 需要 pywin32: pip install pywin32
import shutil
from datetime import timedelta

class PPT2Video:
    def __init__(self, ppt_file, video_file, config_file="config.yaml"):
        self.ppt_file = ppt_file
        self.video_file = video_file
        self.prs = Presentation(ppt_file)
        self.default_duration = 2  # 默认无声视频时长（秒）
        self.resolution = (1920, 1080)  # 1080p 分辨率
        
        # 从 config.yaml 读取音频目录
        with open(config_file, 'r', encoding='utf-8') as f:
            config = yaml.safe_load(f)
        self.audio_dir = config.get('output_dir', 'audio')
        self.temp_dir = os.path.join(self.audio_dir, "temp_slides")  # 临时目录使用 audio_dir/temp_slides

    def export_slides_to_images(self):
        """使用PowerPoint将每页幻灯片导出为图片，确保顺序正确"""
        if os.path.exists(self.temp_dir):
            shutil.rmtree(self.temp_dir)
        os.makedirs(self.temp_dir)

        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(os.path.abspath(self.ppt_file))
        
        presentation.Export(os.path.abspath(self.temp_dir), "PNG")
        presentation.Close()
        powerpoint.Quit()
        
        exported_files = sorted(
            [f for f in os.listdir(self.temp_dir) if f.lower().endswith(".png")],
            key=lambda x: int(x.split('Slide')[1].split('.')[0])
        )
        slide_images = []
        for i, old_file in enumerate(exported_files):
            old_path = os.path.join(self.temp_dir, old_file)
            new_path = os.path.join(self.temp_dir, f"slide_{i}.png")
            os.rename(old_path, new_path)
            print(f"Generated image for slide {i}: {new_path}")
            slide_images.append(new_path)
        
        return slide_images
    
    def str_time(self, seconds):
        td = timedelta(seconds=seconds)
        h = td.seconds//3600
        m = (td.seconds//60)%60
        s = seconds = td.seconds - h*3600 - m*60
        ms = int((seconds % 1) * 1000)
        return f"{h:02d}:{m:02d}:{s:02d},{ms:03d}"

    def generate_srt_time(self, start_time, duration):
        """生成 SRT 文件的时间格式，确保始终为 HH:MM:SS,mmm --> HH:MM:SS,mmm"""
        # 计算开始和结束时间
        start_seconds = start_time
        end_seconds = start_time + duration
        start_str = self.str_time(start_seconds)
        end_str = self.str_time(end_seconds)
        return f"{start_str} --> {end_str}"
    
    def convert(self):
        """将PPT和音频合成为1080p视频，并生成字幕文件"""
        slide_images = self.export_slides_to_images()
        
        if len(slide_images) != len(self.prs.slides):
            print(f"Error: Number of images ({len(slide_images)}) does not match slides ({len(self.prs.slides)})")
            return
        
        clips = []
        srt_entries = []
        total_time = 0  # 用于计算字幕时间
        
        for i, slide in enumerate(self.prs.slides):
            print(f"Processing slide {i}")
            img_file = slide_images[i]
            
            # 第一页特殊处理：2秒无声视频
            if i == 0:
                clip = ImageClip(img_file, duration=self.default_duration).resized(width=self.resolution[0], height=self.resolution[1])
                clips.append(clip)
                print(f"  Created 2-second silent clip for slide {i}, start: {total_time:.2f}s, end: {total_time + self.default_duration:.2f}s")
                total_time += self.default_duration
                continue
            
            # 查找该页的文本和音频文件
            text_files = sorted(
                [f for f in os.listdir(self.audio_dir) if f.startswith(f"slide-{i:03d}-") and f.endswith(".txt")],
                key=lambda x: int(x.split('-')[2].split('.')[0])
            )
            audio_files = sorted(
                [f for f in os.listdir(self.audio_dir) if f.startswith(f"slide-{i:03d}-") and f.endswith(".wav")],
                key=lambda x: int(x.split('-')[2].split('.')[0])
            )
            
            # 检查文本和音频是否匹配
            if len(text_files) != len(audio_files):
                print(f"Error: Mismatch between text files ({len(text_files)}) and audio files ({len(audio_files)}) for slide {i}")
                return
            
            if not text_files:  # 无文本和音频
                clip = ImageClip(img_file, duration=self.default_duration).resized(width=self.resolution[0], height=self.resolution[1])
                clips.append(clip)
                print(f"  Created 2-second silent clip for slide {i}, start: {total_time:.2f}s, end: {total_time + self.default_duration:.2f}s")
                total_time += self.default_duration
            else:  # 有文本和音频
                audio_clips = [AudioFileClip(os.path.join(self.audio_dir, af)) for af in audio_files]
                total_duration = sum(ac.duration for ac in audio_clips)
                combined_audio = concatenate_audioclips(audio_clips)
                
                clip = ImageClip(img_file, duration=total_duration).resized(width=self.resolution[0], height=self.resolution[1])
                clip = clip.with_audio(combined_audio)
                clips.append(clip)
                print(f"  Created clip for slide {i} with audio, total duration {total_duration:.2f} seconds, start: {total_time:.2f}s, end: {total_time + total_duration:.2f}s")
                
                # 生成字幕并打印时间信息
                current_time = total_time
                for j, (txt_file, audio_clip) in enumerate(zip(text_files, audio_clips)):
                    with open(os.path.join(self.audio_dir, txt_file), 'r', encoding='utf-8') as f:
                        text = f.read().strip()
                    srt_entry = f"{len(srt_entries) + 1}\n"
                    srt_entry += f"{self.generate_srt_time(current_time, audio_clip.duration)}\n"
                    srt_entry += f"{text}\n\n"
                    srt_entries.append(srt_entry)
                    print(f"    Added text/audio {txt_file}: start {current_time:.2f}s, end {current_time + audio_clip.duration:.2f}s, duration {audio_clip.duration:.2f}s")
                    current_time += audio_clip.duration
                
                total_time += total_duration
        
        if clips:
            final_video = concatenate_videoclips(clips)
            final_video.write_videofile(
                self.video_file,
                fps=24,
                codec="libx264",
                bitrate="5000k",
                preset="medium",
                audio_codec="aac"
            )
            print(f"Video saved as {self.video_file} (1080p)")
            
            srt_file = os.path.join(self.audio_dir, os.path.splitext(os.path.basename(self.video_file))[0] + ".srt")
            with open(srt_file, 'w', encoding='utf-8') as f:
                f.write("".join(srt_entries))
            print(f"Subtitles saved as {srt_file}")
            
            shutil.rmtree(self.temp_dir)
        else:
            print("No clips to process, video not generated")